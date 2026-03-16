"""
util/parse_qualitative.py

Extracts plain text from qualitative source files and returns a dict of
document texts keyed by document type — ready to be passed to Prompt C.

── File inventory per quarter ─────────────────────────────────────────────────

FY22Q4/                         (no Performance PDF available)
  FinancialStatementFY22Q4.xlsx → handled by parse_financials.py
  TranscriptFY22Q4.docx         → "transcript"
  SlidesFY22Q4.pptx             → "earnings_slides"
  OutlookFY22Q4.pptx            → "outlook_slides"
  PressReleaseFY22Q4.docx       → "press_release"
  FY22Q4ProductList.docx        → "product_releases"

FY2023/Qn/FY23Qn-zip/
  FinancialStatementFY23Q1.xlsx → handled by parse_financials.py
  Metrics_FY23Q1.xlsx           → not used (investor metrics, not qualitative)
  MSFT_FY23Q1_10Q.docx          → not used (SEC filing, too long/legal)
  TranscriptFY23Q1.docx         → "transcript"
  SlidesFY23Q1.pptx             → "earnings_slides"
  OutlookFY23Q1.pptx            → "outlook_slides"
  PressReleaseFY23Q1.docx       → "press_release"
  FY23Q1ProductList.docx        → "product_releases"
  FY23 Q1 - Cash Flows...pdf    → not used (redundant with XLSX)
  FY23 Q1 - Income Statements...pdf → not used (redundant with XLSX)
  FY23 Q1 - Press Releases...pdf    → not used (redundant with PressRelease.docx)

FY2023/Qn/  (one level up from zip folder)
  FY23 Q1 - Performance - Investor Relations - Microsoft.pdf → "performance"

── File matching ──────────────────────────────────────────────────────────────
load_qualitative_docs() uses case-insensitive substring matching on filenames.
Exact naming convention from your files:
  "Transcript"     → DOCX
  "Slides"         → PPTX  (earnings slides — NOT Outlook)
  "Outlook"        → PPTX
  "PressRelease"   → DOCX
  "ProductList"    → DOCX
  "Performance"    → PDF   (lives one level UP from the zip subfolder)
"""

from __future__ import annotations
import re
from pathlib import Path

from docx import Document
from pptx import Presentation


# ── Individual file parsers ────────────────────────────────────────────────────

def parse_docx(path: Path) -> str:
    """Extract all paragraph text from a .docx file."""
    doc = Document(path)
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def parse_pptx(path: Path) -> str:
    """
    Extract all text frame content from a .pptx file.
    Each slide is prefixed with [Slide N] for orientation.
    """
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        if lines:
            slides.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def parse_pdf(path: Path) -> str:
    """
    Extract text from a PDF using PyMuPDF (fitz).
    Install with: pip install pymupdf
    """
    try:
        import fitz  # type: ignore
        doc = fitz.open(str(path))
        return "\n\n".join(page.get_text() for page in doc)
    except ImportError:
        return f"[PDF parsing unavailable — install pymupdf to parse {path.name}]"


# ── Quarter document loader ────────────────────────────────────────────────────

def load_qualitative_docs(quarter_dir: Path, performance_dir: Path | None = None) -> dict[str, str | None]:
    """
    Find and parse all qualitative documents for a quarter.

    Args:
        quarter_dir:     The inner zip folder, e.g. data/input/FY2023/Q1/FY23Q1-zip/
                         (or data/input/FY22Q4/ for the init quarter)
        performance_dir: The folder one level up from the zip folder where the
                         Performance PDF lives, e.g. data/input/FY2023/Q1/
                         Pass None for FY22Q4 (no Performance PDF available).

    Returns:
        {
            "transcript":       str | None,
            "earnings_slides":  str | None,
            "outlook_slides":   str | None,
            "press_release":    str | None,
            "product_releases": str | None,
            "performance":      str | None,
        }
    """
    def _find(directory: Path, pattern: str) -> Path | None:
        """Return the first file in directory whose name matches pattern."""
        if directory is None or not directory.exists():
            return None
        rx = re.compile(pattern, re.IGNORECASE)
        for f in sorted(directory.iterdir()):
            if f.is_file() and rx.search(f.name):
                return f
        return None

    # ── Match files by the naming conventions in your actual file set ─────────
    # "Slides" matches SlidesFY23Q1.pptx but NOT OutlookFY23Q1.pptx
    transcript_path    = _find(quarter_dir, r"transcript")
    slides_path        = _find(quarter_dir, r"^slides|[_-]slides")
    # Fallback: any pptx that isn't Outlook
    if slides_path is None:
        for f in sorted(quarter_dir.iterdir()):
            if f.suffix.lower() in (".pptx", ".ppt") and "outlook" not in f.name.lower():
                slides_path = f
                break
    outlook_path       = _find(quarter_dir, r"outlook")
    press_release_path = _find(quarter_dir, r"pressrelease|press.?release")
    product_list_path  = _find(quarter_dir, r"productlist|product.?list")

    # Performance PDF lives one level up from the zip folder (FY2023 quarters only)
    perf_path = None
    if performance_dir is not None:
        perf_path = _find(performance_dir, r"performance")

    def _docx(p): return parse_docx(p) if p and p.suffix.lower() == ".docx" else None
    def _pptx(p): return parse_pptx(p) if p and p.suffix.lower() in (".pptx", ".ppt") else None
    def _pdf(p):  return parse_pdf(p)  if p and p.suffix.lower() == ".pdf"  else None

    return {
        "transcript":       _docx(transcript_path),
        "earnings_slides":  _pptx(slides_path),
        "outlook_slides":   _pptx(outlook_path),
        "press_release":    _docx(press_release_path),
        "product_releases": _docx(product_list_path),
        "performance":      _pdf(perf_path),
    }